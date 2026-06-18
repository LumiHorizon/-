# -*- coding: utf-8 -*-
# 描述: PPT 生成模板脚本
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x13, 0x21, 0x3C)
NAVY2 = RGBColor(0x1B, 0x2F, 0x52)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
INK = RGBColor(0x1E, 0x2A, 0x44)
MUTED = RGBColor(0x6B, 0x7A, 0x99)
UP = RGBColor(0xD7, 0x26, 0x3D)      # 红涨
DOWN = RGBColor(0x1B, 0x9E, 0x5A)    # 绿跌
GOLD = RGBColor(0xE6, 0xB3, 0x5A)
BORDER = RGBColor(0xD5, 0xDB, 0xE8)
CREAM = RGBColor(0xFB, 0xF3, 0xE4)

FONT = "PingFang SC"

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]


def new_slide(bg):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def rect(s, x, y, w, h, fill, line=None, oval=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.OVAL if oval else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """paras: list of paragraphs; each paragraph is dict
       {runs: [(txt, size, color, bold, italic)], space_after}"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        if p.get("space_after"):
            para.space_after = Pt(p["space_after"])
        for (txt, size, color, bold, italic) in p["runs"]:
            r = para.add_run()
            r.text = txt
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
    return tb


def P(runs, space_after=None):
    return {"runs": runs, "space_after": space_after}


def R(txt, size, color, bold=False, italic=False):
    return (txt, size, color, bold, italic)


# ---------- Slide 1 封面 ----------
s = new_slide(NAVY)
rect(s, 8.55, 0, 1.45, 5.625, NAVY2)
rect(s, 8.4, 0, 0.12, 5.625, GOLD)
text(s, 0.7, 1.0, 6.5, 0.5, [P([R("2026年6月11日 · 周四", 16, ICE)])])
text(s, 0.7, 1.5, 7.0, 1.3, [P([R("盘面总结", 54, WHITE, bold=True)])])
text(s, 0.7, 2.9, 7.0, 0.6, [P([R("缩量企稳 · 情景A兑现", 24, GOLD, bold=True)])])
text(s, 0.7, 3.6, 7.4, 0.5,
     [P([R("两市成交约2.55万亿（较昨缩量）｜沪指微跌守3987｜科创50逆势翻红", 14, ICE)])])
text(s, 0.7, 4.9, 7.0, 0.4,
     [P([R("数据来源：东方财富 · 框架来源：当日早盘蒸馏", 10, MUTED, italic=True)])])

# ---------- Slide 2 指数收盘 ----------
s = new_slide(LIGHT)
text(s, 0.5, 0.35, 9, 0.7, [P([R("指数收盘表现", 32, INK, bold=True)])])

cards = [
    ("上证指数", "3987.01", "-0.16%", DOWN, "失守4000后收敛"),
    ("深证成指", "14851.98", "-0.68%", DOWN, "跌幅明显收窄"),
    ("创业板指", "3811.25", "-1.13%", DOWN, "仍是最弱一环"),
    ("科创50", "1662.44", "+0.62%", UP, "逆势翻红"),
]
cw, gap, x0, y0, ch = 2.1, 0.23, 0.5, 1.25, 2.0
for i, (name, px, chg, c, note) in enumerate(cards):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, ch, WHITE, line=BORDER)
    rect(s, x, y0, cw, 0.08, c)
    text(s, x + 0.15, y0 + 0.22, cw - 0.3, 0.35, [P([R(name, 14, MUTED, bold=True)])])
    text(s, x + 0.15, y0 + 0.62, cw - 0.3, 0.5, [P([R(px, 24, INK, bold=True)])])
    text(s, x + 0.15, y0 + 1.15, cw - 0.3, 0.45, [P([R(chg, 22, c, bold=True)])])
    text(s, x + 0.15, y0 + 1.62, cw - 0.3, 0.3, [P([R(note, 10.5, MUTED)])])

by = 3.7
rect(s, 0.5, by, 9.03, 1.35, NAVY)
text(s, 0.85, by + 0.28, 2.9, 0.8, [P([R("2.55万亿", 40, GOLD, bold=True)])])
text(s, 0.88, by + 1.0, 2.9, 0.3, [P([R("两市合计成交额", 11, ICE)])])
text(s, 4.1, by + 0.3, 5.2, 0.85, [
    P([R("较昨日 2.62 万亿缩量约 3%", 16, WHITE, bold=True)], space_after=4),
    P([R("缩量 = 抛压衰竭信号，对应早盘框架「开盘量能」判断中的积极情景", 12, ICE)]),
])

# ---------- Slide 3 框架验证 ----------
s = new_slide(LIGHT)
text(s, 0.5, 0.35, 9, 0.7, [P([R("早盘框架验证：情景A兑现", 32, INK, bold=True)])])
text(s, 0.5, 1.0, 9, 0.4,
     [P([R("地量出现后，次日唯一关键变量是开盘量能 —— 不预判，只应对", 14, MUTED, italic=True)])])

ay, chh = 1.6, 2.3
rect(s, 0.5, ay, 4.4, chh, NAVY)
rect(s, 0.5, ay, 0.1, chh, GOLD)
text(s, 0.8, ay + 0.2, 3.9, 0.45,
     [P([R("情景A：缩量  ", 19, WHITE, bold=True), R("已兑现", 14, GOLD, bold=True)])])
text(s, 0.8, ay + 0.75, 3.9, 1.4, [
    P([R("• 抛压衰竭，指数获支撑", 13, ICE)], space_after=6),
    P([R("• 沪指仅微跌 0.16%，守住 3987", 13, ICE)], space_after=6),
    P([R("• 盘中修复概率上升 → 科创50 翻红验证", 13, ICE)]),
])

rect(s, 5.15, ay, 4.4, chh, WHITE, line=BORDER)
text(s, 5.45, ay + 0.2, 3.9, 0.45,
     [P([R("情景B：放量下杀  ", 19, MUTED, bold=True), R("未触发", 14, MUTED)])])
text(s, 5.45, ay + 0.75, 3.9, 1.4, [
    P([R("• 若放量下杀 = 资金借外盘出逃", 13, MUTED)], space_after=6),
    P([R("• 反弹观察期顺延", 13, MUTED)], space_after=6),
    P([R("• 今日未出现该情形", 13, MUTED)]),
])

ny = 4.25
rect(s, 0.5, ny, 9.05, 0.85, CREAM)
rect(s, 0.5, ny, 0.1, 0.85, GOLD)
text(s, 0.8, ny + 0.12, 8.5, 0.6, [
    P([R("注意：", 13, INK, bold=True),
       R("企稳 ≠ 右侧。「放量承接」的右侧确认信号尚未出现，参与级别仍以做T为主，不追涨。", 13, INK)]),
], anchor=MSO_ANCHOR.MIDDLE)

# ---------- Slide 4 结构分化 ----------
s = new_slide(LIGHT)
text(s, 0.5, 0.35, 9, 0.7, [P([R("结构分化：科技内部「高切低」", 32, INK, bold=True)])])

ly, lh = 1.3, 3.7
rect(s, 0.5, ly, 4.4, lh, WHITE, line=BORDER)
rect(s, 0.5, ly, 4.4, 0.08, NAVY)
text(s, 0.8, ly + 0.25, 3.9, 0.4, [P([R("指数层面：修复是结构性的", 17, INK, bold=True)])])
text(s, 0.8, ly + 0.8, 1.9, 0.6, [P([R("+0.62%", 32, UP, bold=True)])])
text(s, 0.8, ly + 1.42, 1.95, 0.3, [P([R("科创50（逆势翻红）", 11, MUTED)])])
text(s, 2.95, ly + 0.8, 1.9, 0.6, [P([R("-1.13%", 32, DOWN, bold=True)])])
text(s, 2.95, ly + 1.42, 1.9, 0.3, [P([R("创业板指（仍最弱）", 11, MUTED)])])
text(s, 0.8, ly + 1.95, 3.9, 1.6, [
    P([R("• 科技不是普涨，而是内部切换", 13, INK)], space_after=6),
    P([R("• 资金从高位远期叙事撤出，向低位硬逻辑集中", 13, INK)], space_after=6),
    P([R("• 操作含义：选结构重于赌指数", 13, INK, bold=True)]),
])

rect(s, 5.15, ly, 4.4, lh, NAVY)
rect(s, 5.15, ly, 4.4, 0.08, GOLD)
text(s, 5.45, ly + 0.25, 3.9, 0.4, [P([R("弱市选股：缺货+涨价 > 远期叙事", 17, WHITE, bold=True)])])
text(s, 5.45, ly + 0.8, 2.2, 0.6, [P([R("+200%", 32, GOLD, bold=True)])])
text(s, 5.45, ly + 1.42, 3.9, 0.3, [P([R("六氟化钨同比涨幅（1670–1810元/kg）", 11, ICE)])])
text(s, 5.45, ly + 1.95, 3.9, 1.6, [
    P([R("• 电子化学品 / 电子树脂 / 六氟化钨：当期业绩可兑现", 13, ICE)], space_after=6),
    P([R("• 下半年海外供应缺口仍在，涨价逻辑未破坏", 13, ICE)], space_after=6),
    P([R("• AI泡沫担忧下，硬逻辑比故事更让资金踏实", 13, ICE)]),
])

# ---------- Slide 5 明日关注 ----------
s = new_slide(NAVY)
text(s, 0.5, 0.35, 9, 0.7, [P([R("明日盘前关注", 32, WHITE, bold=True)])])

rows = [
    ("1", "高低切换是否成立",
     "看昨日抗跌品种（MLCC / 铜箔 / 液冷材料端）是否主动领涨；若转为补跌 = 科技进入普遍出清，反弹期拉长"),
    ("2", "修复质量决定参与深度",
     "光模块 / 数据中心电源等超跌方向：放量承接才建仓加仓；缩量弱反抽只做T，不恋战"),
    ("3", "纪律红线",
     "不见量价确认不加仓（右侧原则）；防御资产（红利 / 银行 / 电力）连续大涨后警惕拥挤交易"),
]
y = 1.25
for n, t, d in rows:
    rect(s, 0.5, y, 9.05, 1.1, NAVY2)
    rect(s, 0.78, y + 0.3, 0.5, 0.5, GOLD, oval=True)
    text(s, 0.78, y + 0.3, 0.5, 0.5, [P([R(n, 20, NAVY, bold=True)])],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.55, y + 0.12, 7.8, 0.38, [P([R(t, 17, GOLD, bold=True)])])
    text(s, 1.55, y + 0.52, 7.75, 0.52, [P([R(d, 12.5, ICE)])])
    y += 1.32

text(s, 0.5, 5.18, 9.05, 0.35,
     [P([R("框架沉淀：地量后看开盘量能（不预判只应对）｜弱市硬逻辑优于远期叙事｜缩量反抽 vs 放量承接", 11, MUTED, italic=True)])])

out = "/Users/lihu/Desktop/stock/bilibili/盘面总结_2026-06-11.pptx"
prs.save(out)
print("PPTX written:", out)
